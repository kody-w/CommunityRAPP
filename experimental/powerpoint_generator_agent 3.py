"""
PowerPoint Generator Agent - McKinsey Style
Purpose: Generate professional consulting-style PowerPoint presentations

Design principles (McKinsey style):
- Clean, minimal design with generous white space
- Georgia serif for titles, Arial for body text
- Primary accent: McKinsey blue (#0033A0)
- Data-driven layouts with clear metrics
- Exhibit labels and source citations
- Professional geometric icons (no emojis)
- Thin accent bar at top of slides
- Clear visual hierarchy

Slide types:
- title: Title slide with exhibit label
- content: Content with bullet points
- comparison: Side-by-side comparison table
- pipeline: Process flow with numbered steps
- stats: Key metrics display
- quote: Key insight highlight
- mixed: Multi-section McKinsey layout
- framework: 2x2 or strategic framework

Usage:
1. Simple: action="create_slide", slide_type="stats", title="Key Metrics"
2. Full presentation: action="create_presentation", slides=[{...}, {...}]
3. Quick RAPP slide: action="create_rapp_slide"
"""

import json
import logging
import os
from datetime import datetime
from typing import Dict, Any, List, Optional, Tuple
from agents.basic_agent import BasicAgent
from utils.storage_factory import get_storage_manager

# Import python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.oxml.ns import nsmap
    from lxml import etree
    PPTX_AVAILABLE = True
except ImportError as e:
    PPTX_AVAILABLE = False
    PPTX_IMPORT_ERROR = str(e)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PowerPointGeneratorAgent(BasicAgent):
    """
    Agent for generating McKinsey-style professional presentations.
    Clean, minimal design with data-driven layouts.
    """

    # McKinsey color palette
    COLORS = {
        "mckinsey_blue": "0033A0",
        "dark_blue": "002266",
        "light_blue": "7EB8FF",
        "black": "000000",
        "dark_gray": "333333",
        "medium_gray": "666666",
        "light_gray": "999999",
        "border_gray": "E0E0E0",
        "background_gray": "F8F9FA",
        "white": "FFFFFF",
        "red": "CC0000",
        "green": "006644",
    }

    # Font configurations
    FONTS = {
        "title": {"name": "Georgia", "size": 28, "bold": False},
        "subtitle": {"name": "Arial", "size": 14, "bold": False},
        "section_header": {"name": "Arial", "size": 12, "bold": True},
        "body": {"name": "Arial", "size": 14, "bold": False},
        "metric_value": {"name": "Georgia", "size": 32, "bold": False},
        "metric_label": {"name": "Arial", "size": 11, "bold": False},
        "exhibit_label": {"name": "Arial", "size": 11, "bold": False},
        "source": {"name": "Arial", "size": 10, "bold": False},
        "step_number": {"name": "Arial", "size": 14, "bold": True},
        "step_label": {"name": "Arial", "size": 12, "bold": True},
    }

    # SVG-style icon definitions (created with shapes)
    ICONS = {
        "document": "rectangle",      # Document/transcript
        "code": "chevron",           # Code/development
        "play": "triangle",          # Demo/play
        "check": "checkmark",        # Validation/test
        "arrow_right": "arrow",      # Flow arrow
        "circle": "oval",            # Numbered step
        "gear": "hexagon",           # Settings/process
        "chart": "bar_chart",        # Analytics
        "user": "person",            # User/customer
        "clock": "clock",            # Time
        "target": "target",          # Goal/objective
        "lightbulb": "diamond",      # Insight/idea
    }

    SLIDE_TYPES = {
        "title": "Title slide with exhibit label and key message",
        "content": "Content slide with bullet points",
        "comparison": "Side-by-side comparison table",
        "pipeline": "Process flow with numbered steps",
        "stats": "Key metrics display with boxes",
        "quote": "Key insight highlight box",
        "mixed": "Multi-section McKinsey layout",
        "framework": "2x2 strategic framework",
        "timeline": "Horizontal timeline view"
    }

    def __init__(self):
        self.name = 'PowerPointGenerator'
        self.metadata = {
            "name": self.name,
            "description": """Generate McKinsey-style professional PowerPoint presentations.

Clean, minimal design with:
- Georgia serif titles, Arial body text
- McKinsey blue (#0033A0) accent color
- Data-driven layouts with clear metrics
- Professional geometric icons (no emojis)
- Exhibit labels and source citations

Actions:
- create_presentation: Create multi-slide presentation
- create_slide: Create a single slide
- create_rapp_slide: Quick RAPP overview slide (McKinsey style)
- list_slide_types: List available slide types

Slide types: title, content, comparison, pipeline, stats, quote, mixed, framework, timeline

Example:
{
  "action": "create_slide",
  "slide_type": "comparison",
  "title": "RAPP enables same-day prototyping",
  "left_items": ["Weeks of requirements", "Months to prototype"],
  "right_items": ["Single discovery call", "Working demo same day"],
  "left_label": "Traditional",
  "right_label": "RAPP Approach"
}""",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "description": "Action to perform",
                        "enum": ["create_presentation", "create_slide", "create_rapp_slide", "list_slide_types"]
                    },
                    "slides": {
                        "type": "array",
                        "description": "Array of slide configurations for create_presentation",
                        "items": {"type": "object"}
                    },
                    "slide_type": {
                        "type": "string",
                        "description": "Type of slide to create",
                        "enum": ["title", "content", "comparison", "pipeline", "stats", "quote", "mixed", "framework", "timeline"]
                    },
                    "title": {
                        "type": "string",
                        "description": "Slide title (can include **bold** markers)"
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "Slide subtitle/description"
                    },
                    "exhibit_number": {
                        "type": "integer",
                        "description": "Exhibit number for the slide (e.g., 1, 2, 3)"
                    },
                    "content": {
                        "type": "array",
                        "description": "Array of bullet points",
                        "items": {"type": "string"}
                    },
                    "left_items": {
                        "type": "array",
                        "description": "Left column items for comparison",
                        "items": {"type": "string"}
                    },
                    "right_items": {
                        "type": "array",
                        "description": "Right column items for comparison",
                        "items": {"type": "string"}
                    },
                    "left_label": {
                        "type": "string",
                        "description": "Label for left column"
                    },
                    "right_label": {
                        "type": "string",
                        "description": "Label for right column"
                    },
                    "steps": {
                        "type": "array",
                        "description": "Array of pipeline steps",
                        "items": {
                            "type": "object",
                            "properties": {
                                "label": {"type": "string"},
                                "description": {"type": "string"}
                            }
                        }
                    },
                    "stats": {
                        "type": "array",
                        "description": "Array of metrics",
                        "items": {
                            "type": "object",
                            "properties": {
                                "value": {"type": "string"},
                                "label": {"type": "string"},
                                "unit": {"type": "string"}
                            }
                        }
                    },
                    "quote": {
                        "type": "string",
                        "description": "Quote or key insight text"
                    },
                    "source": {
                        "type": "string",
                        "description": "Source citation for the slide"
                    },
                    "output_filename": {
                        "type": "string",
                        "description": "Output filename (without .pptx)"
                    },
                    "save_to_storage": {
                        "type": "boolean",
                        "description": "Save to Azure storage (default: true)"
                    }
                },
                "required": ["action"]
            }
        }
        super().__init__(self.name, self.metadata)

        try:
            self.storage = get_storage_manager()
        except Exception as e:
            logger.warning(f"Storage not available: {e}")
            self.storage = None

    def perform(self, **kwargs) -> str:
        """Execute the requested action."""
        if not PPTX_AVAILABLE:
            return json.dumps({
                "status": "error",
                "error": f"python-pptx library not available: {PPTX_IMPORT_ERROR}",
                "suggestion": "Install with: pip install python-pptx"
            })

        action = kwargs.get('action', 'create_slide')

        try:
            if action == 'list_slide_types':
                return self._list_slide_types()
            elif action == 'create_rapp_slide':
                return self._create_rapp_slide(**kwargs)
            elif action == 'create_slide':
                return self._create_slide(**kwargs)
            elif action == 'create_presentation':
                return self._create_presentation(**kwargs)
            else:
                return json.dumps({
                    "status": "error",
                    "error": f"Unknown action: {action}",
                    "available_actions": ["create_presentation", "create_slide", "create_rapp_slide", "list_slide_types"]
                })
        except Exception as e:
            logger.error(f"PowerPoint generation error: {e}")
            import traceback
            return json.dumps({
                "status": "error",
                "error": str(e),
                "traceback": traceback.format_exc()
            })

    def _list_slide_types(self) -> str:
        """List available slide types."""
        return json.dumps({
            "status": "success",
            "slide_types": self.SLIDE_TYPES,
            "style": "McKinsey consulting style"
        }, indent=2)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    def _create_base_presentation(self) -> Presentation:
        """Create a base presentation with 16:9 aspect ratio."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        return prs

    def _add_accent_bar(self, slide, prs) -> None:
        """Add McKinsey blue accent bar at top of slide."""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(0.08)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])
        bar.line.fill.background()

    def _add_exhibit_label(self, slide, exhibit_number: int, x: float = 0.5, y: float = 0.25) -> None:
        """Add exhibit label (e.g., 'Exhibit 1')."""
        label_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2), Inches(0.3))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"EXHIBIT {exhibit_number}"
        p.font.name = self.FONTS["exhibit_label"]["name"]
        p.font.size = Pt(self.FONTS["exhibit_label"]["size"])
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])

    def _add_title_with_highlight(self, slide, title: str, x: float, y: float, width: float) -> None:
        """Add title with **bold** text highlighting."""
        title_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        # Parse title for **bold** markers
        import re
        parts = re.split(r'(\*\*[^*]+\*\*)', title)

        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                # Bold text (McKinsey blue)
                run = p.add_run()
                run.text = part[2:-2]
                run.font.name = self.FONTS["title"]["name"]
                run.font.size = Pt(self.FONTS["title"]["size"])
                run.font.bold = True
                run.font.color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])
            elif part:
                # Regular text
                run = p.add_run()
                run.text = part
                run.font.name = self.FONTS["title"]["name"]
                run.font.size = Pt(self.FONTS["title"]["size"])
                run.font.color.rgb = self._hex_to_rgb(self.COLORS["black"])

    def _add_subtitle(self, slide, subtitle: str, x: float, y: float, width: float) -> None:
        """Add subtitle text."""
        sub_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = self.FONTS["subtitle"]["name"]
        p.font.size = Pt(self.FONTS["subtitle"]["size"])
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])

    def _add_section_header(self, slide, text: str, x: float, y: float, width: float) -> None:
        """Add section header with underline."""
        # Header text
        header_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.4))
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = text.upper()
        p.font.name = self.FONTS["section_header"]["name"]
        p.font.size = Pt(self.FONTS["section_header"]["size"])
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])

        # Underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y + 0.35),
            Inches(width * 0.3), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])
        line.line.fill.background()

    def _add_source_citation(self, slide, source: str, page_number: int = 1) -> None:
        """Add source citation and page number at bottom."""
        # Divider line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(6.9),
            Inches(12.333), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        line.line.fill.background()

        # Source text
        source_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(10), Inches(0.3))
        tf = source_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Source: {source}"
        p.font.name = self.FONTS["source"]["name"]
        p.font.size = Pt(self.FONTS["source"]["size"])
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["light_gray"])

        # Page number
        page_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.3))
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(page_number)
        p.font.name = self.FONTS["source"]["name"]
        p.font.size = Pt(self.FONTS["source"]["size"])
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["light_gray"])
        p.alignment = PP_ALIGN.RIGHT

    def _add_divider_line(self, slide, x: float, y1: float, y2: float) -> None:
        """Add vertical divider line."""
        height = y2 - y1
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y1),
            Inches(0.01), Inches(height)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        line.line.fill.background()

    def _add_numbered_circle(self, slide, number: int, x: float, y: float, size: float = 0.45) -> None:
        """Add a numbered circle (McKinsey-style step indicator)."""
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])
        circle.line.fill.background()

        # Number text
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.name = self.FONTS["step_number"]["name"]
        p.font.size = Pt(self.FONTS["step_number"]["size"])
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["white"])
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)

    def _add_bullet_indicator(self, slide, x: float, y: float, color: str = "mckinsey_blue", size: float = 0.1) -> None:
        """Add a small circular bullet indicator."""
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y + 0.05),
            Inches(size), Inches(size)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS[color])
        bullet.line.fill.background()

    def _add_arrow_connector(self, slide, x: float, y: float, width: float = 0.3) -> None:
        """Add an arrow connector between elements."""
        arrow_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.3))
        tf = arrow_box.text_frame
        p = tf.paragraphs[0]
        p.text = "→"
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["light_gray"])
        p.alignment = PP_ALIGN.CENTER

    def _add_metric_box(self, slide, value: str, label: str, x: float, y: float,
                        width: float = 2.5, height: float = 1.5, unit: str = "") -> None:
        """Add a metric display box."""
        # Box border
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["white"])
        box.line.color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        box.line.width = Pt(1)

        # Value
        value_box = slide.shapes.add_textbox(
            Inches(x), Inches(y + 0.3),
            Inches(width), Inches(0.6)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.text = value
        run.font.name = self.FONTS["metric_value"]["name"]
        run.font.size = Pt(self.FONTS["metric_value"]["size"])
        run.font.color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])

        if unit:
            unit_run = p.add_run()
            unit_run.text = unit
            unit_run.font.name = self.FONTS["metric_value"]["name"]
            unit_run.font.size = Pt(16)
            unit_run.font.color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x), Inches(y + height - 0.5),
            Inches(width), Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label.upper()
        p.font.name = self.FONTS["metric_label"]["name"]
        p.font.size = Pt(self.FONTS["metric_label"]["size"])
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])
        p.alignment = PP_ALIGN.CENTER

    def _add_insight_box(self, slide, text: str, x: float, y: float,
                         width: float, height: float) -> None:
        """Add a key insight box (McKinsey blue background)."""
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])
        box.line.fill.background()

        # Parse text for highlights
        import re
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 0.3),
            Inches(width - 0.6), Inches(height - 0.6)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        parts = re.split(r'(\*\*[^*]+\*\*)', text)
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                run = p.add_run()
                run.text = part[2:-2]
                run.font.name = "Georgia"
                run.font.size = Pt(18)
                run.font.italic = True
                run.font.bold = True
                run.font.color.rgb = self._hex_to_rgb(self.COLORS["light_blue"])
            elif part:
                run = p.add_run()
                run.text = part
                run.font.name = "Georgia"
                run.font.size = Pt(18)
                run.font.italic = True
                run.font.color.rgb = self._hex_to_rgb(self.COLORS["white"])

    def _add_pipeline_box(self, slide, steps: List[Dict], x: float, y: float,
                          width: float, height: float) -> None:
        """Add a pipeline flow box with numbered steps."""
        # Background box
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["background_gray"])
        box.line.color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        box.line.width = Pt(1)

        # Calculate step positions
        num_steps = len(steps)
        step_width = (width - 1) / num_steps
        start_x = x + 0.5

        for i, step in enumerate(steps):
            step_x = start_x + (i * step_width)
            step_center_x = step_x + (step_width / 2) - 0.225

            # Numbered circle
            self._add_numbered_circle(slide, i + 1, step_center_x, y + 0.4)

            # Step label
            label_box = slide.shapes.add_textbox(
                Inches(step_x), Inches(y + 1.0),
                Inches(step_width), Inches(0.4)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = step.get('label', f'Step {i+1}')
            p.font.name = self.FONTS["step_label"]["name"]
            p.font.size = Pt(self.FONTS["step_label"]["size"])
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])
            p.alignment = PP_ALIGN.CENTER

            # Step description
            if step.get('description'):
                desc_box = slide.shapes.add_textbox(
                    Inches(step_x), Inches(y + 1.35),
                    Inches(step_width), Inches(0.3)
                )
                tf = desc_box.text_frame
                p = tf.paragraphs[0]
                p.text = step['description']
                p.font.name = "Arial"
                p.font.size = Pt(10)
                p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])
                p.alignment = PP_ALIGN.CENTER

            # Arrow between steps
            if i < num_steps - 1:
                arrow_x = step_x + step_width - 0.2
                self._add_arrow_connector(slide, arrow_x, y + 0.45, 0.4)

    def _add_comparison_table(self, slide, left_items: List[str], right_items: List[str],
                              left_label: str, right_label: str,
                              x: float, y: float, width: float) -> None:
        """Add a comparison table with bullet indicators."""
        # Column labels
        col_width = (width - 0.8) / 2

        # Left label (red indicator)
        left_label_box = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y),
            Inches(col_width), Inches(0.4)
        )
        tf = left_label_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_label.upper()
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["red"])

        # Right label (blue indicator)
        right_label_box = slide.shapes.add_textbox(
            Inches(x + col_width + 0.8), Inches(y),
            Inches(col_width), Inches(0.4)
        )
        tf = right_label_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_label.upper()
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])

        # Rows
        max_items = max(len(left_items), len(right_items))
        row_height = 0.55

        for i in range(max_items):
            row_y = y + 0.5 + (i * row_height)

            # Row divider line
            if i < max_items:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x), Inches(row_y + row_height - 0.05),
                    Inches(width), Inches(0.01)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
                line.line.fill.background()

            # Left item
            if i < len(left_items):
                self._add_bullet_indicator(slide, x, row_y + 0.1, "red")

                left_text = slide.shapes.add_textbox(
                    Inches(x + 0.25), Inches(row_y),
                    Inches(col_width), Inches(row_height)
                )
                tf = left_text.text_frame
                p = tf.paragraphs[0]
                p.text = left_items[i]
                p.font.name = "Arial"
                p.font.size = Pt(14)
                p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])

            # Arrow
            self._add_arrow_connector(slide, x + col_width + 0.1, row_y, 0.5)

            # Right item
            if i < len(right_items):
                self._add_bullet_indicator(slide, x + col_width + 0.7, row_y + 0.1, "mckinsey_blue")

                right_text = slide.shapes.add_textbox(
                    Inches(x + col_width + 0.95), Inches(row_y),
                    Inches(col_width), Inches(row_height)
                )
                tf = right_text.text_frame
                p = tf.paragraphs[0]
                p.text = right_items[i]
                p.font.name = "Arial"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])

    # ==================== SLIDE TYPES ====================

    def _add_title_slide(self, prs: Presentation, config: Dict, page_num: int = 1) -> None:
        """Add a McKinsey-style title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_accent_bar(slide, prs)

        title = config.get('title', 'Presentation Title')
        subtitle = config.get('subtitle', '')
        exhibit = config.get('exhibit_number', page_num)

        # Centered title area
        self._add_exhibit_label(slide, exhibit, 0.5, 2.5)
        self._add_title_with_highlight(slide, title, 0.5, 2.9, 12.333)

        if subtitle:
            self._add_subtitle(slide, subtitle, 0.5, 3.8, 12.333)

        source = config.get('source', 'Internal analysis')
        self._add_source_citation(slide, source, page_num)

    def _add_content_slide(self, prs: Presentation, config: Dict, page_num: int = 1) -> None:
        """Add a content slide with bullet points."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_accent_bar(slide, prs)

        title = config.get('title', 'Content')
        subtitle = config.get('subtitle', '')
        content = config.get('content', [])
        exhibit = config.get('exhibit_number', page_num)

        # Header
        self._add_exhibit_label(slide, exhibit)
        self._add_title_with_highlight(slide, title, 0.5, 0.55, 12.333)

        # Title underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        line.line.fill.background()

        if subtitle:
            self._add_subtitle(slide, subtitle, 0.5, 1.35, 12.333)

        # Content bullets
        y_start = 2.0 if subtitle else 1.6
        for i, item in enumerate(content):
            self._add_bullet_indicator(slide, 0.7, y_start + (i * 0.6) + 0.08, "mckinsey_blue")

            item_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(y_start + (i * 0.6)),
                Inches(11.0), Inches(0.5)
            )
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])

        source = config.get('source', 'Internal analysis')
        self._add_source_citation(slide, source, page_num)

    def _add_comparison_slide(self, prs: Presentation, config: Dict, page_num: int = 1) -> None:
        """Add a comparison slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_accent_bar(slide, prs)

        title = config.get('title', 'Comparison')
        subtitle = config.get('subtitle', '')
        exhibit = config.get('exhibit_number', page_num)

        self._add_exhibit_label(slide, exhibit)
        self._add_title_with_highlight(slide, title, 0.5, 0.55, 12.333)

        # Title underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        line.line.fill.background()

        if subtitle:
            self._add_subtitle(slide, subtitle, 0.5, 1.35, 12.333)

        # Comparison table
        self._add_comparison_table(
            slide,
            config.get('left_items', []),
            config.get('right_items', []),
            config.get('left_label', 'Before'),
            config.get('right_label', 'After'),
            0.5, 2.0 if subtitle else 1.6, 12.333
        )

        source = config.get('source', 'Internal analysis')
        self._add_source_citation(slide, source, page_num)

    def _add_pipeline_slide(self, prs: Presentation, config: Dict, page_num: int = 1) -> None:
        """Add a pipeline/process flow slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_accent_bar(slide, prs)

        title = config.get('title', 'Process Flow')
        subtitle = config.get('subtitle', '')
        exhibit = config.get('exhibit_number', page_num)
        steps = config.get('steps', [
            {"label": "Step 1", "description": "First step"},
            {"label": "Step 2", "description": "Second step"},
            {"label": "Step 3", "description": "Third step"},
            {"label": "Step 4", "description": "Fourth step"}
        ])

        self._add_exhibit_label(slide, exhibit)
        self._add_title_with_highlight(slide, title, 0.5, 0.55, 12.333)

        # Title underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        line.line.fill.background()

        if subtitle:
            self._add_subtitle(slide, subtitle, 0.5, 1.35, 12.333)

        # Pipeline
        self._add_pipeline_box(slide, steps, 0.5, 2.5, 12.333, 2.0)

        source = config.get('source', 'Internal analysis')
        self._add_source_citation(slide, source, page_num)

    def _add_stats_slide(self, prs: Presentation, config: Dict, page_num: int = 1) -> None:
        """Add a statistics/metrics slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_accent_bar(slide, prs)

        title = config.get('title', 'Key Metrics')
        subtitle = config.get('subtitle', '')
        exhibit = config.get('exhibit_number', page_num)
        stats = config.get('stats', [
            {"value": "100", "unit": "%", "label": "Metric One"},
            {"value": "50", "unit": "+", "label": "Metric Two"},
            {"value": "24/7", "unit": "", "label": "Metric Three"}
        ])

        self._add_exhibit_label(slide, exhibit)
        self._add_title_with_highlight(slide, title, 0.5, 0.55, 12.333)

        # Title underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        line.line.fill.background()

        if subtitle:
            self._add_subtitle(slide, subtitle, 0.5, 1.35, 12.333)

        # Metrics boxes
        num_stats = len(stats)
        box_width = min(3.5, 11.0 / num_stats)
        spacing = (12.333 - (box_width * num_stats)) / (num_stats + 1)
        y_start = 2.5

        for i, stat in enumerate(stats):
            x = 0.5 + spacing + (i * (box_width + spacing))
            self._add_metric_box(
                slide,
                stat.get('value', ''),
                stat.get('label', ''),
                x, y_start,
                box_width, 2.0,
                stat.get('unit', '')
            )

        source = config.get('source', 'Internal analysis')
        self._add_source_citation(slide, source, page_num)

    def _add_quote_slide(self, prs: Presentation, config: Dict, page_num: int = 1) -> None:
        """Add a quote/key insight slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_accent_bar(slide, prs)

        quote = config.get('quote', 'Key insight goes here')
        author = config.get('quote_author', '')
        exhibit = config.get('exhibit_number', page_num)

        self._add_exhibit_label(slide, exhibit)

        # Large insight box
        self._add_insight_box(slide, f'"{quote}"', 1.0, 2.0, 11.333, 3.0)

        # Author attribution
        if author:
            author_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.333), Inches(0.5))
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {author}"
            p.font.name = "Arial"
            p.font.size = Pt(14)
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])
            p.alignment = PP_ALIGN.RIGHT

        source = config.get('source', 'Internal analysis')
        self._add_source_citation(slide, source, page_num)

    def _add_mixed_slide(self, prs: Presentation, config: Dict, page_num: int = 1) -> None:
        """Add a mixed McKinsey-style layout slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._add_accent_bar(slide, prs)

        title = config.get('title', 'Overview')
        subtitle = config.get('subtitle', '')
        exhibit = config.get('exhibit_number', page_num)

        self._add_exhibit_label(slide, exhibit)
        self._add_title_with_highlight(slide, title, 0.5, 0.55, 12.333)

        # Title underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        line.line.fill.background()

        if subtitle:
            self._add_subtitle(slide, subtitle, 0.5, 1.35, 12.333)

        source = config.get('source', 'Internal analysis')
        self._add_source_citation(slide, source, page_num)

    # ==================== RAPP SLIDE ====================

    def _create_rapp_slide(self, **kwargs) -> str:
        """Create the full RAPP overview slide (McKinsey style)."""
        output_filename = kwargs.get('output_filename', 'RAPP_Overview_McKinsey')

        prs = self._create_base_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Accent bar
        self._add_accent_bar(slide, prs)

        # Exhibit label
        self._add_exhibit_label(slide, 1)

        # Title with highlight
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        run1 = p.add_run()
        run1.text = "RAPP enables "
        run1.font.name = "Georgia"
        run1.font.size = Pt(28)
        run1.font.color.rgb = self._hex_to_rgb(self.COLORS["black"])

        run2 = p.add_run()
        run2.text = "same-day prototyping"
        run2.font.name = "Georgia"
        run2.font.size = Pt(28)
        run2.font.bold = True
        run2.font.color.rgb = self._hex_to_rgb(self.COLORS["mckinsey_blue"])

        run3 = p.add_run()
        run3.text = ", reducing time-to-demo from months to hours"
        run3.font.name = "Georgia"
        run3.font.size = Pt(28)
        run3.font.color.rgb = self._hex_to_rgb(self.COLORS["black"])

        # Title underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
        line.line.fill.background()

        # Subtitle
        self._add_subtitle(slide, "Rapid Agent Production Pipeline transforms discovery calls into working demonstrations", 0.5, 1.35, 12.333)

        # Vertical divider
        self._add_divider_line(slide, 6.666, 1.8, 6.7)

        # ===== LEFT COLUMN =====

        # Section header
        self._add_section_header(slide, "Traditional vs. RAPP Approach", 0.5, 1.85, 5.8)

        # Comparison items
        comparisons = [
            ("Weeks of requirements gathering", "Single discovery call"),
            ("Months to first prototype", "Working demo same day"),
            ("High risk, late validation", "Early feedback, low risk"),
            ('"Trust us, it will work"', '"Here, try it yourself"'),
        ]

        y_start = 2.5
        for i, (trad, rapp) in enumerate(comparisons):
            row_y = y_start + (i * 0.55)

            # Row divider
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(row_y + 0.5),
                Inches(5.8), Inches(0.01)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["border_gray"])
            line.line.fill.background()

            # Red bullet
            self._add_bullet_indicator(slide, 0.5, row_y + 0.08, "red")

            # Traditional text
            trad_box = slide.shapes.add_textbox(Inches(0.7), Inches(row_y), Inches(2.3), Inches(0.4))
            tf = trad_box.text_frame
            p = tf.paragraphs[0]
            p.text = trad
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])

            # Arrow
            self._add_arrow_connector(slide, 3.0, row_y - 0.05, 0.4)

            # Blue bullet
            self._add_bullet_indicator(slide, 3.5, row_y + 0.08, "mckinsey_blue")

            # RAPP text
            rapp_box = slide.shapes.add_textbox(Inches(3.7), Inches(row_y), Inches(2.5), Inches(0.4))
            tf = rapp_box.text_frame
            p = tf.paragraphs[0]
            p.text = rapp
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])

        # Pipeline box
        steps = [
            {"label": "Transcript", "description": "Discovery call"},
            {"label": "Agent Code", "description": "Auto-generated"},
            {"label": "Demo JSON", "description": "Scenario data"},
            {"label": "HTML Tester", "description": "Live preview"}
        ]
        self._add_pipeline_box(slide, steps, 0.5, 4.8, 5.8, 1.8)

        # ===== RIGHT COLUMN =====

        # Section header
        self._add_section_header(slide, "Speed Advantage", 7.0, 1.85, 5.8)

        # Metrics
        metrics = [
            {"value": "<1", "unit": " day", "label": "Discovery to demo"},
            {"value": "Day 2", "unit": "", "label": "Customer feedback"},
            {"value": "Week 1", "unit": "", "label": "Production ready"}
        ]

        metric_width = 1.7
        metric_spacing = 0.3
        start_x = 7.0

        for i, metric in enumerate(metrics):
            x = start_x + (i * (metric_width + metric_spacing))
            self._add_metric_box(slide, metric["value"], metric["label"], x, 2.5, metric_width, 1.5, metric["unit"])

        # Key insight box
        self._add_insight_box(
            slide,
            '"A **10-minute demo** closes more deals than a **100-page proposal**"',
            7.0, 4.3, 5.833, 1.2
        )

        # Source citation
        self._add_source_citation(slide, "RAPP Pipeline internal metrics; client engagement data 2024-2025", 1)

        return self._save_presentation(prs, output_filename, kwargs)

    def _create_slide(self, **kwargs) -> str:
        """Create a single slide."""
        slide_type = kwargs.get('slide_type', 'content')
        output_filename = kwargs.get('output_filename', f'slide_{slide_type}')

        prs = self._create_base_presentation()

        config = {k: v for k, v in kwargs.items()}

        slide_methods = {
            'title': self._add_title_slide,
            'content': self._add_content_slide,
            'comparison': self._add_comparison_slide,
            'pipeline': self._add_pipeline_slide,
            'stats': self._add_stats_slide,
            'quote': self._add_quote_slide,
            'mixed': self._add_mixed_slide
        }

        method = slide_methods.get(slide_type, self._add_content_slide)
        method(prs, config, 1)

        return self._save_presentation(prs, output_filename, kwargs)

    def _create_presentation(self, **kwargs) -> str:
        """Create a full presentation with multiple slides."""
        slides = kwargs.get('slides', [])
        output_filename = kwargs.get('output_filename', 'presentation')

        if not slides:
            return json.dumps({
                "status": "error",
                "error": "No slides provided. Use 'slides' parameter with array of slide configs."
            })

        prs = self._create_base_presentation()

        slide_methods = {
            'title': self._add_title_slide,
            'content': self._add_content_slide,
            'comparison': self._add_comparison_slide,
            'pipeline': self._add_pipeline_slide,
            'stats': self._add_stats_slide,
            'quote': self._add_quote_slide,
            'mixed': self._add_mixed_slide
        }

        for i, slide_config in enumerate(slides):
            slide_type = slide_config.get('type', 'content')
            method = slide_methods.get(slide_type, self._add_content_slide)
            method(prs, slide_config, i + 1)

        return self._save_presentation(prs, output_filename, kwargs)

    def _save_presentation(self, prs: Presentation, filename: str, kwargs: Dict) -> str:
        """Save the presentation to file."""
        save_to_storage = kwargs.get('save_to_storage', True)
        storage_path = kwargs.get('storage_path', 'presentations')

        if not filename.endswith('.pptx'):
            filename = f"{filename}.pptx"

        local_path = f"/tmp/{filename}"
        prs.save(local_path)

        result = {
            "status": "success",
            "filename": filename,
            "local_path": local_path,
            "style": "McKinsey consulting style"
        }

        if save_to_storage and self.storage:
            try:
                with open(local_path, 'rb') as f:
                    content = f.read()

                storage_file_path = f"{storage_path}/{filename}"
                try:
                    self.storage.write_file('presentations', storage_file_path, content)
                    result["storage_path"] = f"presentations/{storage_file_path}"
                except Exception:
                    result["storage_note"] = "Could not save to Azure storage (binary file)"
            except Exception as e:
                result["storage_error"] = str(e)

        return json.dumps(result, indent=2)
